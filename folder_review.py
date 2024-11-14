import streamlit as st
import os
from docx import Document as DocxDocument
from pptx import Presentation
import pdfplumber

from config import *
from prompt import *
from metrics import *
from groq import Groq

client = Groq(api_key=groq_api_key)  # Replace with your API key

def process_file(file):
    if file.name.endswith(".txt"):
        return file.getvalue().decode("utf-8")
    elif file.name.endswith(".docx"):
        doc = DocxDocument(file)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file.name.endswith(".pdf"):
        with pdfplumber.open(file) as pdf:
            text = ""
            for page in pdf.pages:
                text += page.extract_text() + "\n"
            return text
    elif file.name.endswith(".pptx"):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file.name.endswith((".py", ".java", ".js", ".html", ".css", ".sql", ".cs", ".c", ".cpp")):
        return file.getvalue().decode("utf-8")
    else:
        st.error("Unsupported file format.")
        return ""
def read_code_file(uploaded_file):
    if uploaded_file is not None:
        return uploaded_file.read().decode("utf-8").splitlines()
    return []

# Function to get code review and improved code from LLaMA model
def get_code_review(file_content, file_name):
    

    prompt = (
        f"Review the following code in the file {file_name}. Provide an explanation of the code, "
        f"identify any errors for each type of error, and display the type of error in bold as a title. "
        f"Include all the errors belonging to that type, using numbering in different lines, and include line numbers for each error identified.\n"
        f"Identify any vulnerabilities and potential improvements.\n"
        f"{file_content}\n\n"
        f"Please provide the output in the following format:\n\n"
        f"--- Code Review for {file_name} ---\n"
        f"************* Module {file_name.split('.')[0]}\n"
        f"<review_content_here>\n"
        f"------------------------------------------------------------------\n"
        f"Your code has been rated at <rating>/10\n"
        f"--- Errors ---\n"
    )  
    try:
        # Call the LLaMA model to generate the review and improved code
        prompt_response = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="Llama3-70b-8192",
        )

        # Access the correct attribute for the response content
        if prompt_response and prompt_response.choices:
            review_content = prompt_response.choices[0].message.content
            # Parse the review content to get errors
            errors_by_type = parse_review_content(review_content)
            return review_content, errors_by_type
        else:
            return "No response received from the model.", {}

    except Exception as e:
        if '401' in str(e):
            return "Authentication failed. Please check your API key.", {}
        else:
            return f"Error while fetching code review: {str(e)}", {}

def parse_review_content(review_content):
    # Initialize the dictionary to hold errors by type
    errors_by_type = {
        'Syntax Errors': [],
        'Semantic Errors': [],
        'Compilation Errors': [],
        'Structural Errors': [],
        'Logic Errors': [],
        'Runtime Errors': [],
        'Unused Code': [],
        'Unused Variables': [],
        'Vulnerabilities': [],
        'Improvemets': [],
        'Improvement Suggestions': [],
        'Improved Code': []  # Separate section for improved code
    }

    current_error_type = None
    lines = review_content.strip().split('\n')

    for line in lines:
        line = line.strip()

        # Detecting the error type section headers
        if "Syntax Errors" in line:
            current_error_type = 'Syntax Errors'
        elif "Semantic Errors" in line:
            current_error_type = 'Semantic Errors'
        elif "Structural Errors" in line:
            current_error_type = 'Structural Errors'
        elif "Compilation Errors" in line:
            current_error_type = 'Compilation Errors'
        elif "Logical Errors" in line or "Logic Errors" in line:
            current_error_type = 'Logic Errors'
        elif "Runtime Errors" in line:
            current_error_type = 'Runtime Errors'
        elif "Unused Code" in line:
            current_error_type = 'Unused Code'
        elif "Unused Variables" in line:
            current_error_type = 'Unused Variables'
        elif "Vulnerabilities" in line:
            current_error_type = 'Vulnerabilities'
        else:
            # If in an error section, append the line to the corresponding section
            if current_error_type and line:
                errors_by_type[current_error_type].append(line)

    return errors_by_type


# Assuming SUPPORTED_CODE_TYPES is defined somewhere
SUPPORTED_CODE_TYPES = ['.py', '.js', '.java', '.c', '.cpp', '.html', '.css', '.sql', '.cs']  # Example extensions
def perform_code_review(directory_path, folder_number, org_std, developer_mode):
    org_std_text = process_file(org_std)
    org_std_text = truncate_text(org_std_text)
    if os.path.exists(directory_path):
        code_reviews = {}
        folder_path = os.path.join(directory_path, folder_number)

        if os.path.exists(folder_path):
            for filename in os.listdir(folder_path):
                if any(filename.endswith(ext) for ext in SUPPORTED_CODE_TYPES):
                    file_path = os.path.join(folder_path, filename)

                    with open(file_path, 'r') as file:
                        file_content = file.read()
                    language_map = {
                        ".py": "Python",
                        ".java": "Java",
                        ".js": "JavaScript",
                        ".html": "HTML",
                        ".css": "CSS",
                        ".sql": "SQL",
                        ".cs": ".NET",
                        ".c": "C",
                        ".cpp": "C++"
                    }

                    # Extract the file extension for new code file
                    file_extension = os.path.splitext(filename)[1]

                    # Set the language dynamically for the new code file
                    language = language_map.get(file_extension, "Unknown Language")

                    # Additional metadata placeholders
                    author = "Unknown"  # Placeholder for author's name
                    reviewer = "Unknown"  # Placeholder for reviewer's name
                    # Generate review prompt with truncated text
                    final_output, review, review1, errors_by_type= generate_review_prompt(org_std_text, file_content, changes=None,
                                                                                        new_code_file_name=filename,
                                                                                        old_code_file_name=None,
                                                                                        developer_mode=developer_mode,
                                                                                        language=language, author=author,
                                                                                        reviewer=reviewer)
                    
                    code_reviews[filename] = (review, errors_by_type)

            # Display all code reviews with error types as nested tabs
            if code_reviews:
                tabs = st.tabs(code_reviews.keys())  # Create tabs with file names
                for tab, (filename, (review, errors_by_type)) in zip(tabs, code_reviews.items()):
                    with tab:
                        try:
                            # Display response with the file name in the heading
                            st.subheader(f"--- Code Review for {filename} ---") 
                            # Display the content with filename references replaced
                            review_output = st.empty()
                            review_output.markdown(f"<div class='response-box'>{final_output}</div>", unsafe_allow_html=True)

                            # Calculate score based on the organization standards
                            score_explain = calculate_score(org_std_text, file_content)
                            if developer_mode:
                                st.write(f"Explanation: {score_explain}")  # Display dynamic score based on quality standards
                            
                            errors_by_type = {key: value for key, value in errors_by_type.items() if value}
                            error_tabs_labels = [
                                f"{error_type} ({len(errors)})" 
                                for error_type, errors in errors_by_type.items() 
                                if errors
                            ]
                            error_tabs_labels.append("Improvements")

                            if error_tabs_labels:  # Only create tabs if there are error types with errors
                                error_tabs = st.tabs(error_tabs_labels)
                                for error_tab, (error_type, errors) in zip(error_tabs, errors_by_type.items()):
                                # Only create a tab if there are errors of that type
                                    with error_tab:
                                        st.write(f"### **{error_type}**")
                                        for i, error in enumerate(errors, start=1):
                                            st.write(f"{error}")
                                with error_tabs[len(errors_by_type)]:
                                    st.write("### **Improvements**")
                                    st.write(review)
                            
                            error_count = {f"{error_type}": len(errors) for error_type, errors in errors_by_type.items()}
                            errors = error_count
                            
                            # Display the calculated total severity score and its classification
                            total_score = calculate_severity(error_count)  # Calculate total severity score
                            color, severity = determine_severity_from_score(total_score)  # Determine the severity level

                            # Display the score with color
                            st.markdown(f"### Severity Analysis")
                            st.write(f"Total Errors: {sum(error_count.values())}")
                            total_score_colored = f'<span style="color:{color}; font-size: 20px;">{total_score}</span>'
                            message_colored = f'<span style="color:{color}; font-size: 20px;">{severity}</span>'
                            st.markdown(f'Total Severity Score: {total_score_colored } - Severity: { message_colored }', unsafe_allow_html=True)
                    
                            # st.write(f"Total Severity Score: {total_score_colored }")
                            # st.write(f"Severity: { message_colored }")              
                            vulnerabilities_found = detect_vulnerabilities(file_content)
                            # Display vulnerabilities and improvements
                            if vulnerabilities_found:
                                st.subheader("Detected Vulnerabilities")
                                for vulnerability in vulnerabilities_found:
                                    st.markdown(f"- {vulnerability}")
                            else:
                                st.success("No vulnerabilities detected.")
                            # Determine and display quality level of code
                            # Determine and display quality level of code
                            
                            #st.write(f"Quality Level of Code: {severity}")

                            # Classify URLs in the new code
                            url_classification = classify_urls(file_content)
                            if url_classification:
                                st.subheader("URL Classification")
                                for url, risk in url_classification.items():
                                    st.write(f"{url}: {risk}")                  
                            #review_content = chat_completion.choices[0].message.content.replace("[filename].(extension)", new_code_file_name)
                            #st.write(review_content)
                        except Exception as e:
                            st.error(f"Error: {str(e)}")
            else:
                st.write("No code files found for review.")
        else:
            st.write("Folder number does not exist in the specified directory. Please check the folder number.")
    else:
        st.write("Directory does not exist. Please check the directory path.")
    
