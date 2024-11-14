
import requests
import subprocess
from prompt import *
from metrics import *
from styles import *
from docx import Document as DocxDocument
from pptx import Presentation
import pdfplumber

def process_file(file):
    if file.name.endswith(".txt"):
        return file.getvalue().decode("utf-8")
    elif file.name.endswith(".docx"):
        doc = DocxDocument(file)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file.name.endswith(".pdf"):
        with pdfplumber.open(file) as pdf:
            text = ""
            for page in pdf.pages:
                text += page.extract_text() + "\n"
            return text
    elif file.name.endswith(".pptx"):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file.name.endswith((".py", ".java", ".js", ".html", ".css", ".sql", ".cs", ".c", ".cpp")):
        return file.getvalue().decode("utf-8")
    else:
        st.error("Unsupported file format.")
        return ""

import os


def read_modified_files(modified_files):
    """Read the contents of the modified files.

    Args:
        modified_files (list): List of file paths or dictionaries containing file metadata.

    Returns:
        dict: A dictionary with filenames as keys and file contents as values.
    """
    file_contents = {}

    for file_info in modified_files:
        # Ensure we have the path as a string, not a dictionary or other data structure
        filename = file_info if isinstance(file_info, str) else file_info.get('path')

        if filename and os.path.isfile(filename):
            try:
                with open(filename, 'r', encoding='utf-8') as file:
                    content = file.read().splitlines()  # Read file and split into lines
                    file_contents[filename] = content
            except Exception as e:
                st.write(f"Error reading {filename}: {str(e)}")
        else:
            st.write(f"File {filename} does not exist or is not a valid path.")
    
    return file_contents



def get_modified_files_from_github(repo_url, pr_number, token):
    repo = repo_url.replace("https://github.com/", "")
    url = f"https://api.github.com/repos/{repo}/pulls/{pr_number}/files"
    headers = {"Authorization": f"token {token}"}

    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        files = response.json()
        modified_files = [file for file in files]
        return modified_files
    else:
        st.write(f"Error fetching PR files: {response.status_code} - {response.json()}")
        return []

# Function to fetch pull requests
def get_pull_requests(repo_url, access_token):
    # Extract owner and repo name from URL
    owner_repo = repo_url.replace("https://github.com/", "")
    api_url = f"https://api.github.com/repos/{owner_repo}/pulls"
    headers = {"Authorization": f"token {access_token}"}
    response = requests.get(api_url, headers=headers)
    if response.status_code == 200:
        return response.json()
    else:
        st.error(f"Error: {response.json().get('message', 'Failed to fetch pull requests')}")
        return []

# Function to display file diffs in a pull request
def get_pull_request_diff(repo_url, pr_number, access_token):
    owner_repo = repo_url.replace("https://github.com/", "")
    diff_url = f"https://api.github.com/repos/{owner_repo}/pulls/{pr_number}/files"
    headers = {"Authorization": f"token {access_token}"}
    response = requests.get(diff_url, headers=headers)
    if response.status_code == 200:
        return response.json()
    else:
        st.error(f"Error: {response.json().get('message', 'Failed to fetch file diffs')}")
        return []

# Helper function to style diff content with line numbers and metadata on the second line
def style_diff_content(patch, color_removed="#FFCCCC", color_added="#CCFFCC"):
    lines = patch.splitlines()
    original_content = ""
    modified_content = ""
    orig_line_num = mod_line_num = 1  # Initialize line numbers for original and modified content

    # Identify metadata line (starting with @@) and separate it from code lines
    if lines[0].startswith("@@"):
        metadata = lines[0]
        original_content += f'<div style="font-weight:bold; padding: 4px;">{metadata}</div>'
        modified_content += f'<div style="font-weight:bold; padding: 4px;">{metadata}</div>'
        lines = lines[1:]  # Remove metadata from main lines list

    # Process remaining lines for color-coded changes with line numbers
    for line in lines:
        if line.startswith("-"):
            # Line removed from original content
            original_content += f'<div style="background-color: {color_removed}; padding: 2px;">{orig_line_num}: {line[1:]}</div>'
            orig_line_num += 1
        elif line.startswith("+"):
            # Line added in modified content
            modified_content += f'<div style="background-color: {color_added}; padding: 2px;">{mod_line_num}: {line[1:]}</div>'
            mod_line_num += 1
        else:
            # Unchanged line, increment both counters
            line_content = line[1:] if line.startswith(" ") else line
            original_content += f'<div>{orig_line_num}: {line_content}</div>'
            modified_content += f'<div>{mod_line_num}: {line_content}</div>'
            orig_line_num += 1
            mod_line_num += 1
    
    return original_content, modified_content


def display_pr_review(repo_url,pr_number,org_std,access_token, developer_mode):
  # Assuming this code is within your code review logic
  if org_std:
      org_std_text = process_file(org_std)
      org_std_text = truncate_text(org_std_text)

      modified_files = get_modified_files_from_github(repo_url, pr_number, access_token)
      st.write(modified_files)
      code_reviews = {}

      # Read contents of all modified files
      # Convert the list to a dictionary
      modified_file_contents = {str(index): item for index, item in enumerate(modified_files)}

      print(modified_file_contents)
      # Loop through each modified file content
      for filename, file_content in modified_file_contents.items():
          language_map = {
                        ".py": "Python",
                        ".java": "Java",
                        ".js": "JavaScript",
                        ".html": "HTML",
                        ".css": "CSS",
                        ".sql": "SQL",
                        ".cs": ".NET",
                        ".c": "C",
                        ".cpp": "C++"
                    }

          # Extract the file extension for new code file
          file_extension = os.path.splitext(filename)[1]

          # Set the language dynamically for the new code file
          language = language_map.get(file_extension, "Unknown Language")

          # Additional metadata placeholders
          author = "Unknown"  # Placeholder for author's name
          reviewer = "Unknown"  # Placeholder for reviewer's name
          # Assuming you have the organization standards in org_std
          final_output, review1, review, errors_by_type = generate_review_prompt(org_std_text, file_content, changes=None,
                                                                        new_code_file_name=filename,
                                                                        old_code_file_name=None,
                                                                        developer_mode=developer_mode,
                                                                        language=language, author=author,
                                                                        reviewer=reviewer)

          code_reviews[filename] = (review, errors_by_type)

      # Display all code reviews with error types as nested tabs
      if code_reviews:
          tabs = st.tabs(code_reviews.keys())  # Create tabs with file names
          for tab, (filename, (review, errors_by_type)) in zip(tabs, code_reviews.items()):
              with tab:
                    try:
                      # Display response with the file name in the heading
                        st.subheader(f"--- Code Review for {filename} ---") 
                        # Display the content with filename references replaced
                        review_output = st.empty()
                        review_output.markdown(f"<div class='response-box'>{final_output}</div>", unsafe_allow_html=True)

                        # Calculate score based on the organization standards
                        score_explain = calculate_score(org_std_text, file_content)
                        if developer_mode:
                            st.write(f"Explanation: {score_explain}")  # Display dynamic score based on quality standards
                        
                        errors_by_type = {key: value for key, value in errors_by_type.items() if value}
                        error_tabs_labels = [
                            f"{error_type} ({len(errors)})" 
                            for error_type, errors in errors_by_type.items() 
                            if errors
                        ]
                        error_tabs_labels.append("Improvements")

                        if error_tabs_labels:  # Only create tabs if there are error types with errors
                            error_tabs = st.tabs(error_tabs_labels)
                            for error_tab, (error_type, errors) in zip(error_tabs, errors_by_type.items()):
                            # Only create a tab if there are errors of that type
                                with error_tab:
                                    st.write(f"### **{error_type}**")
                                    for i, error in enumerate(errors, start=1):
                                        st.write(f"{error}")
                            with error_tabs[len(errors_by_type)]:
                                st.write("### **Improvements**")
                                st.write(review)
                        
                        error_count = {f"{error_type}": len(errors) for error_type, errors in errors_by_type.items()}
                        errors = error_count
                        
                        # Display the calculated total severity score and its classification
                        total_score = calculate_severity(error_count)  # Calculate total severity score
                        color, severity = determine_severity_from_score(total_score)  # Determine the severity level

                        # Display the score with color
                        st.markdown(f"### Severity Analysis")
                        st.write(f"Total Errors: {sum(error_count.values())}")
                        total_score_colored = f'<span style="color:{color}; font-size: 20px;">{total_score}</span>'
                        message_colored = f'<span style="color:{color}; font-size: 20px;">{severity}</span>'
                        st.markdown(f'Total Severity Score: {total_score_colored } - Severity: { message_colored }', unsafe_allow_html=True)
                
                        # st.write(f"Total Severity Score: {total_score_colored }")
                        # st.write(f"Severity: { message_colored }")              
                        # vulnerabilities_found = detect_vulnerabilities(file_content)
                        # # Display vulnerabilities and improvements
                        # if vulnerabilities_found:
                        #     st.subheader("Detected Vulnerabilities")
                        #     for vulnerability in vulnerabilities_found:
                        #         st.markdown(f"- {vulnerability}")
                        # else:
                        #     st.success("No vulnerabilities detected.")
                        # Determine and display quality level of code
                        # Determine and display quality level of code
                        
                        st.write(f"Quality Level of Code: {severity}")

                        # Classify URLs in the new code
                        url_classification = classify_urls(file_content)
                        if url_classification:
                            st.subheader("URL Classification")
                            for url, risk in url_classification.items():
                                st.write(f"{url}: {risk}")                  
                        #review_content = chat_completion.choices[0].message.content.replace("[filename].(extension)", new_code_file_name)
                        #st.write(review_content)
                    except Exception as e:
                        st.error(f"Error: {str(e)}")
      else:
          st.write("No code files found for review.")
  else:
      st.write("Please upload the organization standards file and ensure there are modified files.")
